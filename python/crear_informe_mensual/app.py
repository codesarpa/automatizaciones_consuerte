from pptx import Presentation
from pptx.util import Inches, Pt

archivo = 'PLANTILLA INFORME DESARROLLO2.pptx'

#valores de prueba:
MES_ACTUAL = "OCTUBRE 2025"
ID_TICKET = "#208880"
ASUNTO_TICKET = "Solicitud de validación para cambio en la estructura de contraseña de usuarios TAT y CS10"
NOTA_TICKET = "Michael Stiff realiz� mejoras en el robot de servicios publicos de la emsa, debido a diferencias en la pagina web de la emsa con respecto a las respuestas obtenidas del web services, lo que gener� que los pdf's se descargaran da�ados, por lo que se agrega una nueva validaci�n donde se consulta el tipo de contenido que tiene la petici�n a la hora de intentar descargar las facturas."
TICKETS_POR_SLIDE = 3
# contenido = {
#     "%MES_ACTUAL%": MES_ACTUAL,
#     "%ID_TICKET%": ID_TICKET,
#     "%ASUNTO_TICKET%": ASUNTO_TICKET,
#     "%NOTA_TICKET%": NOTA_TICKET
# }
tickets = [
    {
        "%MES_ACTUAL%": "OCTUBRE 2025",
        "%ID_TICKET%": "#208880",
        "%ASUNTO_TICKET%": "MEJORAS ROBOT SERVICIOS PUBLICOS EMSA",
        "%NOTA_TICKET%": "Michael Stiff realizó mejoras en el robot..."
    },
    {
        "%MES_ACTUAL%": "OCTUBRE 2025",
        "%ID_TICKET%": "#208881",
        "%ASUNTO_TICKET%": "OPTIMIZACIÓN ROBOT CONSIGNA",
        "%NOTA_TICKET%": "Se mejoró el flujo de descarga de comprobantes..."
    },
    {
        "%MES_ACTUAL%": "OCTUBRE 2025",
        "%ID_TICKET%": "#208882",
        "%ASUNTO_TICKET%": "ACTUALIZACIÓN VALIDACIÓN DE FACTURAS",
        "%NOTA_TICKET%": "Se agregó control de tipo MIME en descargas..."
    },
    {
        "%MES_ACTUAL%": "OCTUBRE 2025",
        "%ID_TICKET%": "#208883",
        "%ASUNTO_TICKET%": "AJUSTES ROBOT CORREOS",
        "%NOTA_TICKET%": "Se corrigió validación de bandeja de salida..."
    }
]


def generar_presentacion(archivo, tickets, tickets_por_slide):
    ppt = Presentation(archivo)
    slide_base = ppt.slides[1]

    # calculamos cuántas diapositivas necesitamos
    for i in range(0, len(tickets), tickets_por_slide):
        grupo = tickets[i:i + tickets_por_slide]
        slide = ppt.slides.add_slide(slide_base.slide_layout)

        # posición inicial de los cuadros (puedes ajustar)
        top = Inches(1)
        left = Inches(1)
        ancho = Inches(8)
        alto = Inches(1.5)

        for t in grupo:
            # formato del bloque de texto por ticket
            texto_ticket = (
                f"{t['%ID_TICKET%']} - {t['%ASUNTO_TICKET%']}\n"
                f"{t['%NOTA_TICKET%']}\n"
            )
            shape = slide.shapes.add_textbox(left, top, ancho, alto)
            text_frame = shape.text_frame
            text_frame.word_wrap = True

            p = text_frame.add_paragraph()
            p.text = texto_ticket
            p.font.size = Pt(12)
            p.space_after = Pt(6)

            # mover hacia abajo para el siguiente ticket
            top += Inches(2)

    ppt.save("INFORME_TICKETS.pptx")
    print("Presentación generada: INFORME_TICKETS.pptx")

# ─────────────────────────────
generar_presentacion(archivo, tickets, TICKETS_POR_SLIDE)

# def procesar_plantilla(archivo, contenido):
#     ppt = Presentation(archivo)
#     # print(ppt.slides.index)
#     for slide in ppt.slides:
#         for shape in slide.shapes:
#             #LOGICA DE REEMPLAZO DE TEXTO
#             if hasattr(shape, 'text'):
#                 procesar_texto(shape, contenido)
#             else:
#                 print("no es texto")
#     ppt.save(f"{contenido['%MES_ACTUAL%']}.pptx")

# def procesar_texto(shape, contenido):
#     for clave, valor in contenido.items():
#         if clave in shape.text:
#             shape.text = shape.text.replace(clave, valor)
    
# procesar_plantilla(archivo, tickets)